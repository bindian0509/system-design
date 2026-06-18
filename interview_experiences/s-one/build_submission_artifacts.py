from __future__ import annotations

import csv
import re
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DECK_MD = ROOT / "sentinelone-health-center-deck.md"
NARRATIVE_MD = ROOT / "sentinelone-health-center-narrative.md"
ROADMAP_CSV = ROOT / "sentinelone-health-center-roadmap-capacity.csv"


def clean_text(text: str) -> str:
    replacements = {
        "\u2013": "-",
        "\u2014": "-",
        "\u2018": "'",
        "\u2019": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u00a0": " ",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    return text


def parse_deck(md: str) -> list[dict[str, list[str] | str]]:
    parts = [part.strip() for part in re.split(r"\n---\n", md) if part.strip()]
    slides = []
    for part in parts:
        lines = [line.rstrip() for line in part.splitlines()]
        headings = [line for line in lines if line.startswith("# ")]
        title = headings[0].replace("# ", "").strip() if headings else "Slide"
        body = []
        in_code = False
        for line in lines:
            if line.startswith("# "):
                continue
            if line.startswith("```"):
                in_code = not in_code
                continue
            if line.startswith("## "):
                body.append(line.replace("## ", "").strip())
            elif in_code:
                if line.strip():
                    body.append("    " + line)
            elif line.strip():
                body.append(line.strip())
        slides.append({"title": title, "body": body})
    return slides


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=(35, 45, 55)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def make_pptx(slides: list[dict[str, list[str] | str]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg = RGBColor(248, 250, 252)
    dark = (15, 23, 42)
    accent = RGBColor(20, 184, 166)
    muted = (71, 85, 105)

    for idx, slide_info in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.color.rgb = accent

        title = str(slide_info["title"])
        add_textbox(slide, Inches(0.55), Inches(0.35), Inches(11.9), Inches(0.65), title, 28, True, dark)
        add_textbox(slide, Inches(12.25), Inches(0.42), Inches(0.7), Inches(0.35), f"{idx + 1}", 12, False, muted)

        body_lines = list(slide_info["body"])
        y = 1.18
        for line in body_lines:
            if y > 6.75:
                break
            is_section = bool(line and not line.startswith("-") and not line.startswith("    ") and len(line) < 64)
            if is_section:
                add_textbox(slide, Inches(0.72), Inches(y), Inches(11.7), Inches(0.28), line, 15, True, dark)
                y += 0.36
            elif line.startswith("    "):
                add_textbox(slide, Inches(1.0), Inches(y), Inches(10.8), Inches(0.22), line.strip(), 10, False, (51, 65, 85))
                y += 0.23
            else:
                text = line[2:].strip() if line.startswith("- ") else line
                prefix = "• " if line.startswith("- ") else ""
                add_textbox(slide, Inches(0.9), Inches(y), Inches(11.2), Inches(0.34), prefix + text, 14, False, muted)
                y += 0.37

        footer = "SentinelOne Agent Platform | Singularity Health Center"
        add_textbox(slide, Inches(0.72), Inches(7.05), Inches(8.0), Inches(0.2), footer, 8, False, (100, 116, 139))

    prs.save(ROOT / "sentinelone-health-center-deck.pptx")


def make_xlsx() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Roadmap Capacity"

    with ROADMAP_CSV.open(newline="", encoding="utf-8") as f:
        for row in csv.reader(f):
            ws.append(row)

    header_fill = PatternFill("solid", fgColor="0F172A")
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = header_fill
        cell.alignment = Alignment(wrap_text=True, vertical="top")

    for row in ws.iter_rows(min_row=2):
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")

    widths = [14, 24, 44, 24, 30, 42, 30, 34]
    for i, width in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = width

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions
    wb.save(ROOT / "sentinelone-health-center-roadmap-capacity.xlsx")


class SimplePDF:
    def __init__(self, path: Path, pagesize=(612, 792)):
        self.path = path
        self.width, self.height = pagesize
        self.pages: list[list[tuple[int, int, int, str, str]]] = []

    @staticmethod
    def _escape(text: str) -> str:
        return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")

    def add_page(self, lines: list[tuple[int, int, str, str]]):
        page = []
        y = self.height - 54
        for size, leading, font, text in lines:
            if y < 48:
                self.pages.append(page)
                page = []
                y = self.height - 54
            page.append((48, y, size, font, text))
            y -= leading
        self.pages.append(page)

    def save(self):
        objects: list[str] = []
        objects.append("<< /Type /Catalog /Pages 2 0 R >>")
        kids = " ".join(f"{3 + i * 2} 0 R" for i in range(len(self.pages)))
        objects.append(f"<< /Type /Pages /Kids [{kids}] /Count {len(self.pages)} >>")
        for i, page in enumerate(self.pages):
            content_obj = 4 + i * 2
            objects.append(
                f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 {self.width} {self.height}] "
                f"/Resources << /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> "
                f"/F2 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold >> "
                f"/F3 << /Type /Font /Subtype /Type1 /BaseFont /Courier >> >> >> "
                f"/Contents {content_obj} 0 R >>"
            )
            stream_parts = []
            for x, y, size, font, text in page:
                stream_parts.append(f"BT /{font} {size} Tf {x} {y} Td ({self._escape(text)}) Tj ET")
            stream = "\n".join(stream_parts)
            objects.append(f"<< /Length {len(stream.encode('latin-1', errors='replace'))} >>\nstream\n{stream}\nendstream")

        offsets = []
        out = bytearray(b"%PDF-1.4\n")
        for idx, obj in enumerate(objects, 1):
            offsets.append(len(out))
            out.extend(f"{idx} 0 obj\n{obj}\nendobj\n".encode("latin-1", errors="replace"))
        xref = len(out)
        out.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("ascii"))
        for offset in offsets:
            out.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
        out.extend(
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii")
        )
        self.path.write_bytes(out)


def wrap(text: str, width: int) -> list[str]:
    text = clean_text(text).strip()
    if not text:
        return [""]
    words = text.split()
    lines, current = [], ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if len(candidate) <= width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def markdown_to_pdf_lines(md: str, width: int = 88) -> list[tuple[int, int, str, str]]:
    lines: list[tuple[int, int, str, str]] = []
    in_code = False
    for raw in md.splitlines():
        line = clean_text(raw.rstrip())
        if line.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            for part in wrap(line, 92):
                lines.append((8, 10, "F3", part))
            continue
        if line.startswith("# "):
            for part in wrap(line[2:], 48):
                lines.append((18, 24, "F2", part))
            lines.append((8, 8, "F1", ""))
        elif line.startswith("## "):
            lines.append((8, 8, "F1", ""))
            for part in wrap(line[3:], 64):
                lines.append((13, 18, "F2", part))
        elif line.startswith("- "):
            for idx, part in enumerate(wrap(line[2:], width - 4)):
                prefix = "- " if idx == 0 else "  "
                lines.append((10, 14, "F1", prefix + part))
        elif line.strip():
            for part in wrap(line, width):
                lines.append((10, 14, "F1", part))
        else:
            lines.append((8, 8, "F1", ""))
    return lines


def make_pdf(source: Path, target: Path, landscape: bool = False) -> None:
    md = source.read_text(encoding="utf-8")
    if landscape:
        pdf = SimplePDF(target, pagesize=(792, 612))
        slides = parse_deck(md)
        for slide in slides:
            lines: list[tuple[int, int, str, str]] = []
            for part in wrap(str(slide["title"]), 56):
                lines.append((22, 28, "F2", part))
            lines.append((8, 10, "F1", ""))
            for item in slide["body"]:
                item = str(item)
                if item.startswith("    "):
                    for part in wrap(item.strip(), 96):
                        lines.append((8, 10, "F3", part))
                elif item.startswith("- "):
                    for idx, part in enumerate(wrap(item[2:], 88)):
                        lines.append((10, 14, "F1", ("- " if idx == 0 else "  ") + part))
                else:
                    for part in wrap(item, 76):
                        lines.append((12, 17, "F2", part))
            pdf.add_page(lines)
    else:
        pdf = SimplePDF(target)
        pdf.add_page(markdown_to_pdf_lines(md))
    pdf.save()


def main() -> None:
    slides = parse_deck(DECK_MD.read_text(encoding="utf-8"))
    make_pptx(slides)
    make_xlsx()
    make_pdf(DECK_MD, ROOT / "sentinelone-health-center-deck.pdf", landscape=True)
    make_pdf(NARRATIVE_MD, ROOT / "sentinelone-health-center-narrative.pdf")


if __name__ == "__main__":
    main()
