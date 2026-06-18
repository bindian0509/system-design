#!/usr/bin/env python3
"""
Convert a python-pptx PPTX to PDF using reportlab.

This is a purpose-built converter for the executive deck.  It renders shapes
(rectangles, rounded-rects, ovals, text boxes, connectors, tables) with
accurate positioning, colors, and fonts.  It also appends a Notes Appendix.

Usage:
    python3 convert_deck_pdf.py [input.pptx] [output.pdf]
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from reportlab.lib.colors import Color
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import Paragraph, Frame
from reportlab.lib.enums import TA_LEFT


# ──────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────

def emu_to_pt(emu: int) -> float:
    """Convert EMU to points (1 inch = 914400 EMU = 72 pt)."""
    return emu * 72.0 / 914400.0


def rgb_to_color(rgb) -> Color:
    """Convert python-pptx RGBColor to reportlab Color."""
    if rgb is None:
        return Color(0, 0, 0)
    return Color(rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0)


def get_shape_fill_color(shape):
    """Try to get fill color from a shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                return rgb_to_color(fc.rgb)
    except Exception:
        pass
    return None


def get_shape_line_color(shape):
    """Try to get line/border color."""
    try:
        line = shape.line
        if line.color and line.color.rgb:
            return rgb_to_color(line.color.rgb)
    except Exception:
        pass
    return None


def get_text_from_shape(shape) -> list[dict]:
    """Extract text runs from a shape with formatting."""
    runs = []
    try:
        tf = shape.text_frame
    except Exception:
        return runs
    for para in tf.paragraphs:
        for run in para.runs:
            try:
                r = run.font
                size = r.size
                size_pt = size / 12700.0 if size else 10
                bold = r.bold or False
                color = rgb_to_color(r.color.rgb) if r.color and r.color.rgb else Color(0, 0, 0)
            except Exception:
                size_pt = 10
                bold = False
                color = Color(0, 0, 0)
            runs.append({
                "text": run.text,
                "size": min(size_pt, 24),  # cap font size
                "bold": bold,
                "color": color,
            })
        runs.append({"text": "\n", "size": 10, "bold": False, "color": Color(0, 0, 0)})
    return runs


# ──────────────────────────────────────────────────────────────
# Renderer
# ──────────────────────────────────────────────────────────────

class DeckRenderer:
    def __init__(self, pptx_path: str, pdf_path: str):
        self.prs = Presentation(pptx_path)
        self.pdf_path = pdf_path
        # Slide dimensions in points
        self.slide_w = emu_to_pt(self.prs.slide_width)
        self.slide_h = emu_to_pt(self.prs.slide_height)
        self.page_size = (self.slide_w, self.slide_h)
        self.c = canvas.Canvas(pdf_path, pagesize=self.page_size)
        self.c.setTitle("Singularity Health Center – Executive Slide Deck")
        self.c.setAuthor("Bharat")
        self.c.setSubject("Senior Engineering Manager – Agent Platform")

    def render(self):
        for i, slide in enumerate(self.prs.slides):
            self._render_slide(slide, i)
            self.c.showPage()

        # Add notes appendix
        self._render_notes_appendix()

        self.c.save()
        print(f"[OK] PDF saved → {self.pdf_path}")

    def _render_slide(self, slide, slide_idx: int):
        c = self.c
        sh = self.slide_h

        # Background
        bg_color = Color(248/255, 250/255, 252/255)  # default bg
        try:
            bg = slide.background.fill
            if bg.type is not None and bg.fore_color and bg.fore_color.rgb:
                bg_color = rgb_to_color(bg.fore_color.rgb)
        except Exception:
            pass
        c.setFillColor(bg_color)
        c.rect(0, 0, self.slide_w, sh, stroke=0, fill=1)

        # Render each shape
        for shape in slide.shapes:
            self._render_shape(shape, sh)

    def _render_shape(self, shape, page_h: float):
        """Render a single shape."""
        c = self.c

        # Position in reportlab coords (bottom-left origin)
        x = emu_to_pt(shape.left) if shape.left else 0
        y_top = emu_to_pt(shape.top) if shape.top else 0
        w = emu_to_pt(shape.width) if shape.width else 0
        h = emu_to_pt(shape.height) if shape.height else 0
        # Convert to bottom-left origin
        y = page_h - y_top - h

        # Handle tables
        if shape.has_table:
            self._render_table(shape, page_h)
            return

        # Handle connectors (lines/arrows)
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            self._render_connector(shape, page_h)
            return

        # Fill
        fill_color = get_shape_fill_color(shape)
        line_color = get_shape_line_color(shape)

        if fill_color and w > 0 and h > 0:
            c.setFillColor(fill_color)
            if line_color:
                c.setStrokeColor(line_color)
                c.setLineWidth(0.5)
            else:
                c.setStrokeColor(fill_color)
                c.setLineWidth(0)

            # Determine shape type
            shape_name = ""
            try:
                shape_name = shape.auto_shape_type.name if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type else ""
            except Exception:
                pass

            if "OVAL" in shape_name:
                cx = x + w / 2
                cy = y + h / 2
                c.ellipse(x, y, x + w, y + h, stroke=1, fill=1)
            elif "ROUNDED" in shape_name:
                radius = min(w, h) * 0.15
                c.roundRect(x, y, w, h, radius, stroke=1, fill=1)
            else:
                c.rect(x, y, w, h, stroke=0, fill=1)

        # Text
        if shape.has_text_frame:
            self._render_text(shape, x, y, w, h)

    def _render_text(self, shape, x: float, y: float, w: float, h: float):
        """Render text from a shape."""
        c = self.c
        try:
            tf = shape.text_frame
        except Exception:
            return

        margin_l = emu_to_pt(tf.margin_left) if tf.margin_left else 2
        margin_t = emu_to_pt(tf.margin_top) if tf.margin_top else 2
        margin_r = emu_to_pt(tf.margin_right) if tf.margin_right else 2

        text_x = x + margin_l
        text_w = w - margin_l - margin_r
        if text_w <= 0:
            text_w = w

        # Calculate starting y position (top of text area)
        text_y = y + h - margin_t

        for para in tf.paragraphs:
            full_text = ""
            font_size = 10
            font_bold = False
            font_color = Color(0, 0, 0)

            for run in para.runs:
                full_text += run.text
                try:
                    if run.font.size:
                        font_size = min(run.font.size / 12700.0, 24)
                    if run.font.bold:
                        font_bold = True
                    if run.font.color and run.font.color.rgb:
                        font_color = rgb_to_color(run.font.color.rgb)
                except Exception:
                    pass

            if not full_text.strip():
                text_y -= font_size + 2
                continue

            c.setFillColor(font_color)
            font_name = "Helvetica-Bold" if font_bold else "Helvetica"
            c.setFont(font_name, font_size)

            # Simple text wrapping
            lines = self._wrap_text(full_text, font_name, font_size, text_w)
            for line in lines:
                if text_y < y:
                    break
                c.drawString(text_x, text_y - font_size, line)
                text_y -= font_size + 2

    def _wrap_text(self, text: str, font_name: str, font_size: float,
                   max_width: float) -> list[str]:
        """Simple word-wrap."""
        from reportlab.pdfbase.pdfmetrics import stringWidth
        lines = []
        for raw_line in text.split("\n"):
            words = raw_line.split()
            if not words:
                lines.append("")
                continue
            current = words[0]
            for word in words[1:]:
                test = current + " " + word
                if stringWidth(test, font_name, font_size) <= max_width:
                    current = test
                else:
                    lines.append(current)
                    current = word
            lines.append(current)
        return lines

    def _render_connector(self, shape, page_h: float):
        """Render a line/connector."""
        c = self.c
        try:
            x1 = emu_to_pt(shape.begin_x) if shape.begin_x else emu_to_pt(shape.left)
            y1 = page_h - emu_to_pt(shape.begin_y) if shape.begin_y else page_h - emu_to_pt(shape.top)
            x2 = emu_to_pt(shape.end_x) if shape.end_x else emu_to_pt(shape.left + shape.width)
            y2 = page_h - emu_to_pt(shape.end_y) if shape.end_y else page_h - emu_to_pt(shape.top + shape.height)
        except Exception:
            x1 = emu_to_pt(shape.left)
            y1 = page_h - emu_to_pt(shape.top)
            x2 = x1 + emu_to_pt(shape.width)
            y2 = y1 - emu_to_pt(shape.height)

        color = get_shape_line_color(shape) or Color(0.2, 0.25, 0.33)
        try:
            lw = shape.line.width
            lw_pt = lw / 12700.0 if lw else 1.0
        except Exception:
            lw_pt = 1.0

        c.setStrokeColor(color)
        c.setLineWidth(lw_pt)
        c.line(x1, y1, x2, y2)

        # Draw arrowhead
        import math
        dx = x2 - x1
        dy = y2 - y1
        length = math.sqrt(dx*dx + dy*dy)
        if length > 0:
            ux, uy = dx/length, dy/length
            arrow_len = 5
            ax = x2 - arrow_len * ux + arrow_len * 0.4 * uy
            ay = y2 - arrow_len * uy - arrow_len * 0.4 * ux
            bx = x2 - arrow_len * ux - arrow_len * 0.4 * uy
            by = y2 - arrow_len * uy + arrow_len * 0.4 * ux
            c.setFillColor(color)
            path = c.beginPath()
            path.moveTo(x2, y2)
            path.lineTo(ax, ay)
            path.lineTo(bx, by)
            path.close()
            c.drawPath(path, fill=1, stroke=0)

    def _render_table(self, shape, page_h: float):
        """Render a table."""
        c = self.c
        table = shape.table
        tbl_x = emu_to_pt(shape.left)
        tbl_y_top = emu_to_pt(shape.top)
        tbl_w = emu_to_pt(shape.width)
        tbl_h = emu_to_pt(shape.height)

        n_rows = len(table.rows)
        n_cols = len(table.columns)

        # Column widths
        col_widths = [emu_to_pt(table.columns[j].width) for j in range(n_cols)]
        row_height = tbl_h / n_rows

        for r in range(n_rows):
            for col in range(n_cols):
                cell = table.cell(r, col)
                cx = tbl_x + sum(col_widths[:col])
                cy_top = tbl_y_top + r * row_height
                cw = col_widths[col]
                ch = row_height
                cy = page_h - cy_top - ch

                # Cell fill
                try:
                    if cell.fill and cell.fill.type is not None:
                        fc = cell.fill.fore_color
                        if fc and fc.rgb:
                            fill = rgb_to_color(fc.rgb)
                            c.setFillColor(fill)
                            c.rect(cx, cy, cw, ch, stroke=0, fill=1)
                except Exception:
                    pass

                # Cell border
                c.setStrokeColor(Color(0.85, 0.85, 0.85))
                c.setLineWidth(0.3)
                c.rect(cx, cy, cw, ch, stroke=1, fill=0)

                # Cell text
                text = cell.text.strip()
                if text:
                    font_size = 8
                    font_name = "Helvetica"
                    font_color = Color(0.2, 0.25, 0.33)
                    try:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    font_size = min(run.font.size / 12700.0, 14)
                                if run.font.bold:
                                    font_name = "Helvetica-Bold"
                                if run.font.color and run.font.color.rgb:
                                    font_color = rgb_to_color(run.font.color.rgb)
                                break
                            break
                    except Exception:
                        pass

                    c.setFillColor(font_color)
                    c.setFont(font_name, font_size)
                    c.drawString(cx + 4, cy + ch / 2 - font_size / 2, text[:60])

    def _render_notes_appendix(self):
        """Add a Notes Appendix section to the PDF."""
        c = self.c
        styles = getSampleStyleSheet()
        note_style = ParagraphStyle(
            'NoteStyle',
            parent=styles['Normal'],
            fontSize=8,
            leading=10,
            fontName='Helvetica',
            textColor=Color(0.2, 0.25, 0.33),
        )
        heading_style = ParagraphStyle(
            'HeadingStyle',
            parent=styles['Heading2'],
            fontSize=14,
            leading=18,
            fontName='Helvetica-Bold',
            textColor=Color(0.06, 0.09, 0.16),
        )
        appendix_title_style = ParagraphStyle(
            'AppendixTitle',
            parent=styles['Heading1'],
            fontSize=20,
            leading=24,
            fontName='Helvetica-Bold',
            textColor=Color(0.06, 0.09, 0.16),
        )

        # Appendix title page
        c.setFillColor(Color(248/255, 250/255, 252/255))
        c.rect(0, 0, self.slide_w, self.slide_h, stroke=0, fill=1)
        c.setFillColor(Color(0.06, 0.09, 0.16))
        c.setFont("Helvetica-Bold", 28)
        c.drawString(72, self.slide_h - 120, "Presenter Notes Appendix")
        c.setFont("Helvetica", 12)
        c.setFillColor(Color(0.28, 0.33, 0.41))
        c.drawString(72, self.slide_h - 150,
                     "Detailed talking points, anticipated Q&A, and transition cues for each slide.")
        c.showPage()

        # One or more pages per slide's notes
        for i, slide in enumerate(self.prs.slides, 1):
            try:
                ns = slide.notes_slide
                notes_text = ns.notes_text_frame.text if ns else ""
            except Exception:
                notes_text = ""

            if not notes_text.strip():
                continue

            # Render notes as text pages
            c.setFillColor(Color(248/255, 250/255, 252/255))
            c.rect(0, 0, self.slide_w, self.slide_h, stroke=0, fill=1)

            # Header
            c.setFillColor(Color(0.08, 0.72, 0.65))
            c.setFont("Helvetica-Bold", 11)
            c.drawString(54, self.slide_h - 36, f"SLIDE {i:02}")

            c.setFillColor(Color(0.06, 0.09, 0.16))
            c.setFont("Helvetica-Bold", 16)
            # Extract title from notes
            lines = notes_text.strip().split("\n")
            title_line = lines[0] if lines else f"Slide {i}"
            c.drawString(100, self.slide_h - 38, title_line[:80])

            # Separator line
            c.setStrokeColor(Color(0.8, 0.84, 0.88))
            c.setLineWidth(0.8)
            c.line(54, self.slide_h - 48, self.slide_w - 54, self.slide_h - 48)

            # Notes body
            body_text = "\n".join(lines[1:]) if len(lines) > 1 else notes_text
            # Render with basic text
            c.setFillColor(Color(0.2, 0.25, 0.33))
            c.setFont("Helvetica", 9)
            text_y = self.slide_h - 68
            for line in body_text.split("\n"):
                if text_y < 40:
                    c.showPage()
                    c.setFillColor(Color(248/255, 250/255, 252/255))
                    c.rect(0, 0, self.slide_w, self.slide_h, stroke=0, fill=1)
                    text_y = self.slide_h - 50
                    c.setFillColor(Color(0.2, 0.25, 0.33))
                    c.setFont("Helvetica", 9)

                # Bold lines (section headers)
                stripped = line.strip()
                if stripped.startswith("=") or not stripped:
                    text_y -= 4
                    continue
                if stripped.isupper() or stripped.endswith(")") and stripped.startswith(("KEY", "ANTICIPATED", "TRANSITION", "PURPOSE", "WALK", "CLOSING", "IMPORTANT", "OPERATING", "WHY", "THIS", "LATENCY", "DATA", "GATE", "ORG", "SEV", "OBSERV", "EXAMPLE", "STATE", "COMPLEXITY", "BUT", "RISK", "EXPECTED", "TEAM", "GATEWAY", "LEGACY", "LEARNING")):
                    c.setFont("Helvetica-Bold", 9.5)
                    c.setFillColor(Color(0.06, 0.09, 0.16))
                    c.drawString(54, text_y, stripped)
                    c.setFont("Helvetica", 9)
                    c.setFillColor(Color(0.2, 0.25, 0.33))
                    text_y -= 13
                elif stripped.startswith("•") or stripped.startswith("-") or stripped.startswith("→"):
                    c.drawString(66, text_y, stripped)
                    text_y -= 12
                else:
                    # Wrap long lines
                    from reportlab.pdfbase.pdfmetrics import stringWidth
                    max_w = self.slide_w - 120
                    words = stripped.split()
                    if not words:
                        text_y -= 6
                        continue
                    current = words[0]
                    for word in words[1:]:
                        test = current + " " + word
                        if stringWidth(test, "Helvetica", 9) <= max_w:
                            current = test
                        else:
                            c.drawString(54, text_y, current)
                            text_y -= 12
                            current = word
                    c.drawString(54, text_y, current)
                    text_y -= 12

            c.showPage()


# ──────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────

def main():
    root = Path(__file__).resolve().parent
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else str(root / "sentinelone-health-center-mature-deck.pptx")
    pdf_path  = sys.argv[2] if len(sys.argv) > 2 else str(root / "sentinelone-health-center-mature-deck.pdf")

    renderer = DeckRenderer(pptx_path, pdf_path)
    renderer.render()


if __name__ == "__main__":
    main()
