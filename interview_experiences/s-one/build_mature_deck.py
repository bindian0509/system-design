from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "sentinelone-health-center-mature-deck.pptx"


W = Inches(13.333)
H = Inches(7.5)

COLORS = {
    "ink": RGBColor(15, 23, 42),
    "muted": RGBColor(71, 85, 105),
    "line": RGBColor(203, 213, 225),
    "bg": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "teal": RGBColor(20, 184, 166),
    "blue": RGBColor(37, 99, 235),
    "violet": RGBColor(124, 58, 237),
    "amber": RGBColor(217, 119, 6),
    "rose": RGBColor(225, 29, 72),
    "green": RGBColor(22, 163, 74),
    "slate": RGBColor(51, 65, 85),
    "soft_teal": RGBColor(204, 251, 241),
    "soft_blue": RGBColor(219, 234, 254),
    "soft_violet": RGBColor(237, 233, 254),
    "soft_amber": RGBColor(254, 243, 199),
    "soft_rose": RGBColor(255, 228, 230),
    "soft_green": RGBColor(220, 252, 231),
}


def set_slide_bg(slide, color=COLORS["bg"]):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    color = color or COLORS["ink"]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, n, heading, sub=None):
    text_box(slide, 0.45, 0.25, 0.55, 0.32, f"{n:02}", 11, True, COLORS["teal"])
    text_box(slide, 0.95, 0.18, 11.6, 0.45, heading, 24, True, COLORS["ink"])
    if sub:
        text_box(slide, 0.97, 0.65, 11.2, 0.32, sub, 10.5, False, COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.94), Inches(12.45), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.color.rgb = COLORS["line"]


def pill(slide, x, y, w, h, text, fill, font=COLORS["ink"], size=10.5, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = font
    return shape


def card(slide, x, y, w, h, heading, body, accent=COLORS["teal"], fill=COLORS["panel"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    text_box(slide, x + 0.18, y + 0.13, w - 0.3, 0.28, heading, 12, True, COLORS["ink"])
    text_box(slide, x + 0.18, y + 0.5, w - 0.3, h - 0.6, body, 9.3, False, COLORS["muted"])
    return shape


def arrow(slide, x1, y1, x2, y2, color=COLORS["slate"], width=1.5):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def box(slide, x, y, w, h, text, fill, border=None, font=COLORS["ink"], size=10.5, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border or fill
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = font
    return shape


def footer(slide):
    text_box(slide, 0.55, 7.1, 5.4, 0.2, "Singularity Health Center | Agent Platform", 7.5, False, COLORS["muted"])


def add_bullets(slide, x, y, w, h, bullets, size=10.2):
    box_shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["muted"]
        p.space_after = Pt(4)
    return box_shape


def slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 1, "Singularity Health Center", "Senior Engineering Manager plan for a hyperscale Agent Platform initiative")
    text_box(s, 0.75, 1.35, 7.4, 0.55, "Deliver customer-trust outcomes while building a reusable real-time health platform.", 24, True)
    text_box(
        s,
        0.78,
        2.05,
        6.7,
        0.6,
        "The assignment is not only an architecture exercise. It tests whether we can deliver a strategic platform while absorbing live operational load, dependency delays, and senior-engineer disagreement.",
        12,
        False,
        COLORS["muted"],
    )
    card(s, 0.8, 3.05, 2.8, 1.15, "Customer Outcome", "Trusted visibility into agent health across millions of endpoints.", COLORS["teal"])
    card(s, 3.9, 3.05, 2.8, 1.15, "Platform Outcome", "Reusable telemetry, coalescing, alert lifecycle, and replay capabilities.", COLORS["blue"])
    card(s, 7.0, 3.05, 2.8, 1.15, "Execution Outcome", "MVP in controlled rollout, GA after scale, correctness, and migration gates.", COLORS["violet"])
    card(s, 10.1, 3.05, 2.4, 1.15, "Org Outcome", "Team stays focused without ignoring current customer pain.", COLORS["amber"])
    pill(s, 0.85, 4.75, 2.15, 0.48, "<5 min detection", COLORS["soft_teal"])
    pill(s, 3.2, 4.75, 2.15, 0.48, "<200 ms API p95", COLORS["soft_blue"])
    pill(s, 5.55, 4.75, 2.35, 0.48, "Billions events/day", COLORS["soft_violet"])
    pill(s, 8.1, 4.75, 2.1, 0.48, "Replay + backfill", COLORS["soft_green"])
    pill(s, 10.4, 4.75, 2.0, 0.48, "Tenant rollout", COLORS["soft_amber"])
    text_box(s, 0.85, 5.75, 11.6, 0.55, "Leadership thesis: protect customers now, reduce legacy drag, and land the new platform through staged scope and evidence-based decisions.", 16, True, COLORS["ink"])
    footer(s)


def slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 2, "Scope, SLOs, And Delivery Guardrails", "Define what matters before debating technology choices")
    card(s, 0.65, 1.25, 3.8, 1.25, "MVP Signals", "Offline/connectivity loss\nAgent disabled\nAnti-tamper disabled\nLow disk/resource risk", COLORS["teal"])
    card(s, 4.75, 1.25, 3.8, 1.25, "GA Expansion", "Alert lifecycle\nSuppression and drill-down\nShadow migration from heartbeat\nTenant-wide rollout", COLORS["blue"])
    card(s, 8.85, 1.25, 3.8, 1.25, "Explicitly Deferred", "Generic rule builder\nAdvanced analytics\nCustom notifications\nFull legacy retirement", COLORS["amber"])

    headers = ["SLO", "Target", "Why It Matters"]
    rows = [
        ["Detection freshness", "99% < 5 min", "Trustworthy health state"],
        ["Dashboard latency", "p95 < 200 ms", "Operator workflow speed"],
        ["Silent gap detection", "< 5 min to page", "No stale green dashboards"],
        ["Accepted event loss", "0 after durable bus", "Replayable correctness"],
    ]
    table = s.shapes.add_table(5, 3, Inches(0.72), Inches(3.05), Inches(11.9), Inches(2.25)).table
    widths = [2.45, 2.0, 7.45]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["ink"]
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["panel"] if r % 2 else RGBColor(241, 245, 249)
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(9.5)
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["slate"]

    text_box(s, 0.85, 5.85, 11.5, 0.45, "Guardrail: every MVP feature must improve customer actionability or de-risk the platform launch.", 14, True, COLORS["ink"])
    footer(s)


def slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 3, "Reference Architecture", "Cloud-neutral, replayable, and intentionally separated into hot, read, and cold paths")

    stages = [
        ("Agents", COLORS["soft_teal"], 0.45),
        ("Ingestion\nGateway", COLORS["soft_blue"], 1.9),
        ("Durable\nEvent Bus", COLORS["soft_violet"], 3.55),
        ("Stream\nProcessing", COLORS["soft_green"], 5.2),
        ("Rules +\nCoalescing", COLORS["soft_amber"], 6.85),
        ("Alert\nState Store", COLORS["soft_teal"], 8.5),
        ("Read Model\n+ Search", COLORS["soft_blue"], 10.15),
        ("APIs\n+ Console", COLORS["soft_violet"], 11.8),
    ]
    y = 2.0
    for i, (label, fill, x) in enumerate(stages):
        box(s, x, y, 1.25, 0.85, label, fill, COLORS["line"], size=9.6)
        if i < len(stages) - 1:
            arrow(s, x + 1.28, y + 0.42, stages[i + 1][2] - 0.05, y + 0.42, COLORS["slate"], 1.2)

    box(s, 3.55, 4.2, 1.25, 0.72, "Data Lake\nAudit/Replay", COLORS["soft_green"], COLORS["line"], size=9.2)
    arrow(s, 4.18, 2.88, 4.18, 4.15, COLORS["green"], 1.2)
    box(s, 8.5, 4.2, 1.25, 0.72, "Lifecycle\nStore", COLORS["soft_amber"], COLORS["line"], size=9.2)
    arrow(s, 9.12, 2.88, 9.12, 4.15, COLORS["amber"], 1.2)

    card(s, 0.75, 5.45, 3.5, 0.95, "Hot Path", "Event-time detection, coalescing, alert transitions, freshness SLO.", COLORS["teal"])
    card(s, 4.9, 5.45, 3.5, 0.95, "Read Path", "Precomputed dashboard views, bounded queries, freshness metadata.", COLORS["blue"])
    card(s, 9.05, 5.45, 3.5, 0.95, "Cold Path", "Long-term retention, backfill, analytics, and rule tuning.", COLORS["green"])
    footer(s)


def slide_4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 4, "Read Path And Latency Budget", "The UI should query health state, not raw telemetry")

    box(s, 0.8, 1.45, 2.0, 0.75, "Console\nDashboard", COLORS["soft_violet"], COLORS["line"])
    box(s, 3.35, 1.45, 2.0, 0.75, "Health APIs\nBounded Queries", COLORS["soft_blue"], COLORS["line"])
    box(s, 5.9, 1.45, 2.0, 0.75, "Read Model\nAggregates", COLORS["soft_teal"], COLORS["line"])
    box(s, 8.45, 1.45, 2.0, 0.75, "Search Index\nFilters", COLORS["soft_green"], COLORS["line"])
    box(s, 10.95, 1.45, 1.6, 0.75, "Alert\nState", COLORS["soft_amber"], COLORS["line"])
    for x1, x2 in [(2.8, 3.35), (5.35, 5.9), (7.9, 8.45), (10.45, 10.95)]:
        arrow(s, x1, 1.83, x2, 1.83)

    headers = ["Budget Segment", "Target", "Design Choice"]
    rows = [
        ["Auth/gateway", "20-40 ms", "Tenant-scoped claims and request shaping"],
        ["API orchestration", "40-60 ms", "No fanout over raw telemetry"],
        ["Read/search", "70-90 ms", "Precomputed views and bounded filters"],
        ["Serialization/network", "20-30 ms", "Slim response models and pagination"],
    ]
    table = s.shapes.add_table(5, 3, Inches(0.8), Inches(3.0), Inches(7.4), Inches(2.3)).table
    for i, w in enumerate([2.1, 1.45, 3.85]):
        table.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["ink"]
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True
        run.font.size = Pt(9)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["panel"] if r % 2 else RGBColor(241, 245, 249)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(8.6)
            run.font.color.rgb = COLORS["slate"]

    card(s, 8.75, 3.05, 3.75, 0.9, "Graceful Degradation", "If search lags, serve current alert state with explicit freshness metadata.", COLORS["amber"])
    card(s, 8.75, 4.15, 3.75, 0.9, "Query Discipline", "Pagination, capped time windows, no tenant-wide scans from UI paths.", COLORS["blue"])
    card(s, 8.75, 5.25, 3.75, 0.9, "Correctness Model", "At-least-once ingestion plus idempotent alert transitions.", COLORS["green"])
    footer(s)


def slide_5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 5, "Alert Coalescing Design", "Turn noisy telemetry into actionable health state within the 5-minute SLO")

    text_box(s, 0.75, 1.25, 4.2, 0.28, "Example: 50 anti-tamper disabled events over 10 minutes", 12.5, True)
    # Timeline
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(2.05), Inches(5.7), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = COLORS["line"]; line.line.color.rgb = COLORS["line"]
    for i, x in enumerate([0.95, 1.4, 1.8, 2.3, 2.9, 3.45, 4.0, 4.65, 5.2, 5.8, 6.25]):
        fill = COLORS["rose"] if i in [0, 1, 2, 3, 4, 5] else COLORS["amber"]
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.93), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = fill; dot.line.color.rgb = fill
    pill(s, 0.85, 2.42, 1.2, 0.34, "t=0", COLORS["soft_rose"], size=9)
    pill(s, 5.55, 2.42, 1.1, 0.34, "t=10m", COLORS["soft_amber"], size=9)

    box(s, 7.0, 1.2, 1.55, 0.64, "Key By\nTenant+Agent", COLORS["soft_blue"], COLORS["line"], size=9)
    box(s, 9.0, 1.2, 1.55, 0.64, "Window\nState", COLORS["soft_teal"], COLORS["line"], size=9)
    box(s, 11.0, 1.2, 1.55, 0.64, "Alert\nTransition", COLORS["soft_green"], COLORS["line"], size=9)
    arrow(s, 8.55, 1.52, 9.0, 1.52)
    arrow(s, 10.55, 1.52, 11.0, 1.52)

    card(s, 0.85, 3.35, 3.6, 1.05, "State Machine", "Open -> update evidence count -> suppress duplicate -> resolve when healthy signal appears.", COLORS["teal"])
    card(s, 4.75, 3.35, 3.6, 1.05, "Auditability", "Store raw event references and rule version, not every raw event in the UI path.", COLORS["blue"])
    card(s, 8.65, 3.35, 3.6, 1.05, "Decision", "Use stream processing for coalescing; keep rule scope narrow for MVP.", COLORS["green"])

    add_bullets(
        s,
        0.95,
        5.05,
        11.2,
        0.95,
        [
            "Why not cron-first: late detection, expensive scans, duplicate intermediate state, and weak replay semantics.",
            "How to control complexity: fixed MVP rules, explicit windows, idempotent updates, per-rule metrics, and rollback flags.",
        ],
        11,
    )
    footer(s)


def slide_6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 6, "Reliability Model And Sev-1 Runbook", "Design for silent failure detection, not just service uptime")

    stages = [
        ("Ingress\nrate", 0.8, COLORS["soft_blue"]),
        ("Bus\nlag", 2.5, COLORS["soft_violet"]),
        ("Processor\nlag", 4.2, COLORS["soft_teal"]),
        ("Store\nwrites", 5.9, COLORS["soft_green"]),
        ("Index\nfreshness", 7.6, COLORS["soft_amber"]),
        ("Synthetic\nprobe", 9.3, COLORS["soft_rose"]),
        ("API\np95", 11.0, COLORS["soft_blue"]),
    ]
    for i, (label, x, fill) in enumerate(stages):
        box(s, x, 1.35, 1.15, 0.72, label, fill, COLORS["line"], size=8.8)
        if i < len(stages) - 1:
            arrow(s, x + 1.15, 1.71, stages[i + 1][1] - 0.05, 1.71)

    headers = ["Incident Step", "SEM / IC Focus", "Technical Focus"]
    rows = [
        ["Declare", "Single IC, comms owner, customer impact clock", "Freeze risky deploys"],
        ["Scope", "Tenant, region, event type, time window", "Compare stage metrics"],
        ["Mitigate", "Choose customer-safe degraded mode", "Failover, pause rules, restart consumers"],
        ["Recover", "Track backfill completion", "Replay from durable offsets"],
        ["Prevent", "Postmortem and owners", "Synthetic gaps, gates, runbooks"],
    ]
    table = s.shapes.add_table(6, 3, Inches(0.75), Inches(2.75), Inches(11.85), Inches(2.8)).table
    for i, w in enumerate([2.0, 4.35, 5.5]):
        table.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = COLORS["ink"]
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = RGBColor(255, 255, 255)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = COLORS["panel"] if r % 2 else RGBColor(241, 245, 249)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(8.4); run.font.color.rgb = COLORS["slate"]

    text_box(s, 0.9, 6.05, 11.2, 0.35, "Operating principle: a stale green Health Center is worse than an explicit degraded state.", 14, True, COLORS["rose"])
    footer(s)


def slide_7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 7, "Roadmap: MVP, GA, Platformization", "Multi-quarter plan with customer-safety gates")

    quarters = [("Q1 MVP", 2.55), ("Q2 GA", 5.55), ("Q3 Platform", 8.55), ("Q4 Reduce Legacy", 11.1)]
    for label, x in quarters:
        pill(s, x, 1.25, 1.55, 0.38, label, COLORS["ink"], RGBColor(255, 255, 255), 9.5)

    lanes = [
        ("Telemetry + Processing", COLORS["teal"], ["Contract\nAdapter\nTop 4 rules", "Replay\nScale test\nIdempotency", "Reusable\ncoalescing\nframework", "Broader\nhealth catalog"]),
        ("Console + APIs", COLORS["blue"], ["Dashboard\nAlert list\nFeature flag", "Lifecycle\nSearch\nDrill-down", "Trends\nReports\nNotifications", "Legacy UI\nconsolidation"]),
        ("Operations", COLORS["green"], ["SLOs\nSynthetic\nRunbooks", "Game days\nCanaries\nGA gates", "Cost model\nOwnership\nOnboarding", "Decom\nrunbooks"]),
        ("Legacy Offline", COLORS["amber"], ["Stabilize\nfalse offline", "Shadow\ncompare", "Tenant\nmigration", "Retire\nsafe paths"]),
    ]
    y = 1.85
    for lane, accent, items in lanes:
        text_box(s, 0.65, y + 0.14, 1.65, 0.35, lane, 9.2, True, COLORS["ink"])
        for i, item in enumerate(items):
            box(s, 2.55 + i * 3.0, y, 1.9 if i < 3 else 1.65, 0.7, item, RGBColor(255, 255, 255), accent, size=8.1)
            if i < 3:
                arrow(s, 4.45 + i * 3.0, y + 0.35, 5.0 + i * 3.0, y + 0.35, COLORS["line"], 1)
        y += 1.05

    card(s, 0.75, 6.25, 3.8, 0.65, "MVP Gate", "End-to-end synthetic telemetry, top four signals, SLO dashboards, allowlist rollout.", COLORS["teal"])
    card(s, 4.8, 6.25, 3.8, 0.65, "GA Gate", "Scale, replay, correctness, shadow validation, and incident readiness.", COLORS["blue"])
    card(s, 8.85, 6.25, 3.8, 0.65, "Migration Gate", "Measured false-positive/negative rate and rollback path per tenant cohort.", COLORS["amber"])
    footer(s)


def slide_8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 8, "Dependency Delay And Operational Drain", "Absorb real-world disruption without pretending capacity is infinite")

    card(s, 0.75, 1.25, 5.8, 1.0, "Gateway Delayed By Two Months", "Freeze telemetry contract, build adapter boundary, generate synthetic/replayed streams, negotiate thin-slice routing for highest-value signals.", COLORS["violet"])
    card(s, 6.95, 1.25, 5.6, 1.0, "Legacy Offline Escalations +40%", "Treat as customer trust risk; stabilize with a short-lived lane, not a full-team roadmap derailment.", COLORS["rose"])

    # Allocation stacked bar
    text_box(s, 0.85, 2.85, 5.0, 0.28, "10-person team allocation during escalation", 12.5, True)
    segments = [
        ("Health Center build", 6, COLORS["blue"]),
        ("Legacy stabilization", 2, COLORS["rose"]),
        ("QA/release", 1, COLORS["green"]),
        ("Tech lead coordination", 1, COLORS["amber"]),
    ]
    x = 0.85
    total_w = 7.2
    for label, count, color in segments:
        w = total_w * count / 10
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.35), Inches(w), Inches(0.55))
        rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.line.color.rgb = RGBColor(255, 255, 255)
        text_box(s, x + 0.05, 3.48, w - 0.1, 0.18, str(count), 9, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
        text_box(s, x, 4.05, w, 0.38, label, 8.2, False, COLORS["muted"], PP_ALIGN.CENTER)
        x += w

    card(s, 8.45, 2.85, 4.0, 0.85, "Exit Criteria", "Legacy lane ends after top root cause patch, instrumentation, and support-volume trend review.", COLORS["green"])
    card(s, 8.45, 3.95, 4.0, 0.85, "Escalation Rule", "Increase allocation only for active Sev-1/Sev-2 impact across major tenants.", COLORS["amber"])
    card(s, 8.45, 5.05, 4.0, 0.85, "Learning Loop", "Each false-offline root cause becomes a Health Center test, rule refinement, or migration guardrail.", COLORS["blue"])
    footer(s)


def slide_9(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 9, "Technical Conflict Resolution", "Use evidence, decision criteria, and ownership to align senior engineers")

    box(s, 0.85, 1.35, 2.25, 0.78, "Engineer A\nStream processing", COLORS["soft_teal"], COLORS["teal"], size=9.3)
    box(s, 10.25, 1.35, 2.25, 0.78, "Engineer B\nDB + cron", COLORS["soft_amber"], COLORS["amber"], size=9.3)
    box(s, 4.95, 1.25, 3.2, 0.98, "SEM Role\nTurn disagreement into a decision system", COLORS["soft_blue"], COLORS["blue"], size=10)
    arrow(s, 3.1, 1.74, 4.9, 1.74)
    arrow(s, 10.25, 1.74, 8.2, 1.74)

    steps = [
        ("1 Criteria", "Freshness, scale, cost, operability, delivery risk"),
        ("2 Spike", "One week with realistic cardinality and failure cases"),
        ("3 Decision", "ADR with assumptions, rollback, and revisit trigger"),
        ("4 Commit", "Both engineers own part of the winning design"),
    ]
    x = 0.9
    for i, (h, b) in enumerate(steps):
        card(s, x + i * 3.05, 3.0, 2.55, 1.1, h, b, [COLORS["teal"], COLORS["blue"], COLORS["violet"], COLORS["green"]][i])
        if i < 3:
            arrow(s, x + i * 3.05 + 2.55, 3.55, x + (i + 1) * 3.05 - 0.12, 3.55, COLORS["line"], 1.2)

    text_box(s, 0.9, 5.05, 11.5, 0.45, "Expected decision: stream processing for coalescing, constrained to MVP rules and strong operability controls.", 14, True, COLORS["ink"])
    add_bullets(
        s,
        1.05,
        5.65,
        10.8,
        0.8,
        [
            "Engineer A leads the processing design and SLO instrumentation.",
            "Engineer B leads simplicity controls: bounded scope, failure-mode review, cost model, and rollback plan.",
        ],
        10.8,
    )
    footer(s)


def slide_10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    title(s, 10, "Major Risks, Mitigations, And Outcomes", "Close with operating maturity, not a technology shopping list")

    # Risk matrix
    text_box(s, 0.75, 1.22, 4.6, 0.28, "Risk Matrix", 12.5, True)
    x0, y0, cell = 0.8, 1.65, 0.82
    for i in range(4):
        for j in range(4):
            color = [COLORS["soft_green"], COLORS["soft_amber"], COLORS["soft_rose"], COLORS["soft_rose"]][min(i + j, 3)]
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + i * cell), Inches(y0 + (3 - j) * cell), Inches(cell), Inches(cell))
            rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.line.color.rgb = RGBColor(255, 255, 255)
    text_box(s, 0.8, 5.0, 3.3, 0.25, "Likelihood ->", 8.5, False, COLORS["muted"], PP_ALIGN.CENTER)
    text_box(s, 0.35, 2.8, 0.35, 0.5, "Impact", 8.5, False, COLORS["muted"])
    risks = [
        ("G", 3, 3, COLORS["violet"]),
        ("S", 2, 3, COLORS["blue"]),
        ("L", 3, 2, COLORS["rose"]),
        ("Q", 2, 2, COLORS["amber"]),
        ("N", 1, 2, COLORS["green"]),
    ]
    for label, ix, iy, color in risks:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x0 + ix * cell + 0.22), Inches(y0 + (3 - iy) * cell + 0.2), Inches(0.36), Inches(0.36))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.color.rgb = color
        text_box(s, x0 + ix * cell + 0.27, y0 + (3 - iy) * cell + 0.27, 0.2, 0.12, label, 7, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)

    card(s, 4.7, 1.35, 3.7, 0.75, "G Gateway Delay", "Adapter contract, synthetic streams, thin-slice integration.", COLORS["violet"])
    card(s, 8.75, 1.35, 3.7, 0.75, "S Stream Complexity", "Narrow MVP rules, SLOs, runbooks, rollback flags.", COLORS["blue"])
    card(s, 4.7, 2.35, 3.7, 0.75, "L Legacy Drain", "Two-engineer stabilization lane with exit criteria.", COLORS["rose"])
    card(s, 8.75, 2.35, 3.7, 0.75, "Q Query Latency", "Precomputed read models, bounded APIs, freshness metadata.", COLORS["amber"])
    card(s, 4.7, 3.35, 3.7, 0.75, "N Noise / Alert Fatigue", "Coalescing, suppression, tenant rollout, support feedback loop.", COLORS["green"])

    text_box(s, 4.75, 4.75, 7.7, 0.35, "Expected business outcomes", 13, True, COLORS["ink"])
    add_bullets(
        s,
        4.9,
        5.22,
        7.4,
        1.05,
        [
            "More trustworthy agent-health visibility for customers and support.",
            "Reduced false-offline escalations and less legacy operational drag.",
            "Reusable Agent Platform foundation for future health and alerting use cases.",
        ],
        11,
    )
    footer(s)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for fn in [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8, slide_9, slide_10]:
        fn(prs)
    prs.save(OUT)


if __name__ == "__main__":
    build()
