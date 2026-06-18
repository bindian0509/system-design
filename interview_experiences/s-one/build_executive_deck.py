#!/usr/bin/env python3
"""
Build an executive-grade Singularity Health Center slide deck (PPTX + PDF).

Features over the original build_mature_deck.py:
  - Unicode icon glyphs on every slide (shields, clocks, gears, charts, etc.)
  - Richer architecture diagrams with labelled data-flow arrows
  - Timeline / swim-lane roadmap with milestone diamonds
  - Org-chart / team-allocation visual
  - 4×4 risk-matrix with plotted risk items and legend
  - Detailed **presenter notes** on all 10 slides
  - Footer with page number
  - Auto-conversion to PDF via LibreOffice (if available)
"""
from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ──────────────────────────────────────────────────────────────
# Paths
# ──────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent
PPTX_OUT = ROOT / "sentinelone-health-center-mature-deck.pptx"
PDF_OUT  = ROOT / "sentinelone-health-center-mature-deck.pdf"

# ──────────────────────────────────────────────────────────────
# Dimensions & palette
# ──────────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

C = {
    "ink":         RGBColor(15, 23, 42),
    "muted":       RGBColor(71, 85, 105),
    "line":        RGBColor(203, 213, 225),
    "bg":          RGBColor(248, 250, 252),
    "panel":       RGBColor(255, 255, 255),
    "white":       RGBColor(255, 255, 255),
    "teal":        RGBColor(20, 184, 166),
    "blue":        RGBColor(37, 99, 235),
    "violet":      RGBColor(124, 58, 237),
    "amber":       RGBColor(217, 119, 6),
    "rose":        RGBColor(225, 29, 72),
    "green":       RGBColor(22, 163, 74),
    "slate":       RGBColor(51, 65, 85),
    "soft_teal":   RGBColor(204, 251, 241),
    "soft_blue":   RGBColor(219, 234, 254),
    "soft_violet": RGBColor(237, 233, 254),
    "soft_amber":  RGBColor(254, 243, 199),
    "soft_rose":   RGBColor(255, 228, 230),
    "soft_green":  RGBColor(220, 252, 231),
    "dark_blue":   RGBColor(30, 58, 138),
    "alt_row":     RGBColor(241, 245, 249),
}

# Unicode icons used as lightweight iconography
ICONS = {
    "shield":    "\U0001F6E1",   # 🛡
    "chart":     "\U0001F4CA",   # 📊
    "gear":      "\u2699",       # ⚙
    "clock":     "\u23F1",       # ⏱
    "bolt":      "\u26A1",       # ⚡
    "check":     "\u2705",       # ✅
    "warning":   "\u26A0",       # ⚠
    "fire":      "\U0001F525",   # 🔥
    "target":    "\U0001F3AF",   # 🎯
    "rocket":    "\U0001F680",   # 🚀
    "people":    "\U0001F465",   # 👥
    "lock":      "\U0001F512",   # 🔒
    "link":      "\U0001F517",   # 🔗
    "star":      "\u2B50",       # ⭐
    "diamond":   "\u25C6",       # ◆
    "arrow_r":   "\u27A4",       # ➤
    "circle":    "\u25CF",       # ●
    "plant":     "\U0001F331",   # 🌱
    "trophy":    "\U0001F3C6",   # 🏆
    "bar":       "\u2503",       # ┃
}

# ──────────────────────────────────────────────────────────────
# Helper primitives
# ──────────────────────────────────────────────────────────────

def set_bg(slide, color=C["bg"]):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _tb(slide, x, y, w, h, text, size=14, bold=False, color=None,
         align=PP_ALIGN.LEFT, font_name="Calibri"):
    color = color or C["ink"]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top  = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return box


def _title(slide, n, heading, sub=None):
    _tb(slide, 0.45, 0.25, 0.55, 0.32, f"{n:02}", 11, True, C["teal"])
    _tb(slide, 0.95, 0.18, 11.6, 0.45, heading, 24, True, C["ink"])
    if sub:
        _tb(slide, 0.97, 0.65, 11.2, 0.32, sub, 10.5, False, C["muted"])
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.45), Inches(0.94),
                                Inches(12.45), Inches(0.015))
    ln.fill.solid(); ln.fill.fore_color.rgb = C["line"]
    ln.line.color.rgb = C["line"]


def _footer(slide, page_num: int):
    _tb(slide, 0.55, 7.1, 5.4, 0.2,
        "Singularity Health Center  |  Agent Platform", 7.5, False, C["muted"])
    _tb(slide, 12.0, 7.1, 0.8, 0.2,
        str(page_num), 7.5, False, C["muted"], PP_ALIGN.RIGHT)


def _pill(slide, x, y, w, h, text, fill, font=None, size=10.5, bold=True):
    font = font or C["ink"]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame; tf.clear()
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = font
    return shape


def _card(slide, x, y, w, h, heading, body, accent=C["teal"], fill=C["panel"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C["line"]; shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.color.rgb = accent
    _tb(slide, x + 0.18, y + 0.13, w - 0.3, 0.28, heading, 12, True, C["ink"])
    _tb(slide, x + 0.18, y + 0.5,  w - 0.3, h - 0.6, body, 9.3, False, C["muted"])
    return shape


def _arrow(slide, x1, y1, x2, y2, color=None, width=1.5):
    color = color or C["slate"]
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    return conn


def _box(slide, x, y, w, h, text, fill, border=None, font=None, size=10.5, bold=True):
    font = font or C["ink"]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border or fill
    tf = shape.text_frame; tf.clear()
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = font
    return shape


def _bullets(slide, x, y, w, h, items, size=10.2):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.level = 0
        p.font.size = Pt(size); p.font.color.rgb = C["muted"]
        p.space_after = Pt(4)
    return box


def _notes(slide, text: str):
    """Add presenter notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _table(slide, x, y, w, h, headers, rows, col_widths=None):
    """Add a styled table and return the table object."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                        Inches(x), Inches(y),
                                        Inches(w), Inches(h))
    table = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for c, hdr in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = C["ink"]
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = C["white"]; run.font.bold = True; run.font.size = Pt(9.5)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C["panel"] if r % 2 else C["alt_row"]
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(9); run.font.color.rgb = C["slate"]
    return table


# ──────────────────────────────────────────────────────────────
# SLIDES
# ──────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title / Executive Framing"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 1, f"{ICONS['shield']}  Singularity Health Center",
           "Senior Engineering Manager plan for a hyperscale Agent Platform initiative")

    _tb(s, 0.75, 1.35, 7.4, 0.55,
        "Deliver customer-trust outcomes while building a reusable real-time health platform.",
        24, True)
    _tb(s, 0.78, 2.05, 6.7, 0.6,
        "The assignment is not only an architecture exercise. It tests whether we can "
        "deliver a strategic platform while absorbing live operational load, dependency "
        "delays, and senior-engineer disagreement.", 12, False, C["muted"])

    # Outcome cards with icons
    _card(s, 0.8,  3.05, 2.8, 1.15,
          f"{ICONS['target']}  Customer Outcome",
          "Trusted visibility into agent health across millions of endpoints.", C["teal"])
    _card(s, 3.9,  3.05, 2.8, 1.15,
          f"{ICONS['gear']}  Platform Outcome",
          "Reusable telemetry, coalescing, alert lifecycle, and replay capabilities.", C["blue"])
    _card(s, 7.0,  3.05, 2.8, 1.15,
          f"{ICONS['rocket']}  Execution Outcome",
          "MVP in controlled rollout, GA after scale, correctness, and migration gates.", C["violet"])
    _card(s, 10.1, 3.05, 2.4, 1.15,
          f"{ICONS['people']}  Org Outcome",
          "Team stays focused without ignoring current customer pain.", C["amber"])

    # KPI pills with icons
    _pill(s, 0.85,  4.75, 2.15, 0.48, f"{ICONS['clock']}  <5 min detection", C["soft_teal"])
    _pill(s, 3.2,   4.75, 2.15, 0.48, f"{ICONS['bolt']}  <200 ms API p95",   C["soft_blue"])
    _pill(s, 5.55,  4.75, 2.35, 0.48, f"{ICONS['chart']}  Billions events/day", C["soft_violet"])
    _pill(s, 8.1,   4.75, 2.1,  0.48, f"{ICONS['link']}  Replay + backfill",  C["soft_green"])
    _pill(s, 10.4,  4.75, 2.0,  0.48, f"{ICONS['lock']}  Tenant rollout",     C["soft_amber"])

    _tb(s, 0.85, 5.75, 11.6, 0.55,
        "Leadership thesis: protect customers now, reduce legacy drag, and land the new "
        "platform through staged scope and evidence-based decisions.", 16, True, C["ink"])

    _footer(s, 1)

    _notes(s, """\
PRESENTER NOTES — Slide 1: Executive Framing
=============================================

OPENING (60 seconds)
• Start by anchoring on the CUSTOMER problem, not the technology.
  "Customers today cannot reliably tell whether their SentinelOne agents are healthy."
• Frame the assignment as a leadership test, not just architecture.

KEY TALKING POINTS
1. Four outcomes — Customer / Platform / Execution / Org — show you understand
   this is more than a feature build.
2. KPIs are concrete and measurable. Call out the <5-minute detection freshness
   as the hardest constraint that drives the architecture.
3. Leadership thesis: staged scope, evidence-based decisions, protect current
   customers. This separates you from candidates who only talk technology.

ANTICIPATED QUESTIONS
• "Why not start with all anomaly types?"
  → Explain ruthless MVP scoping; we prove correctness on 4 signals first.
• "How do you balance new-platform velocity with legacy maintenance?"
  → Preview the dedicated stabilization lane (Slide 8).

TRANSITION → "Let me show you how I scoped the MVP and set SLOs before
discussing architecture choices."
""")


def slide_02_scope(prs):
    """Scope, SLOs, And Delivery Guardrails"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 2, f"{ICONS['target']}  Scope, SLOs, And Delivery Guardrails",
           "Define what matters before debating technology choices")

    _card(s, 0.65, 1.25, 3.8, 1.25,
          f"{ICONS['check']}  MVP Signals",
          "Offline / connectivity loss\nAgent disabled\n"
          "Anti-tamper disabled\nLow disk / resource risk", C["teal"])
    _card(s, 4.75, 1.25, 3.8, 1.25,
          f"{ICONS['rocket']}  GA Expansion",
          "Alert lifecycle\nSuppression and drill-down\n"
          "Shadow migration from heartbeat\nTenant-wide rollout", C["blue"])
    _card(s, 8.85, 1.25, 3.8, 1.25,
          f"{ICONS['warning']}  Explicitly Deferred",
          "Generic rule builder\nAdvanced analytics\n"
          "Custom notifications\nFull legacy retirement", C["amber"])

    _table(s, 0.72, 3.05, 11.9, 2.25,
           ["SLO", "Target", "Why It Matters"],
           [
               ["Detection freshness",   "99% < 5 min",             "Trustworthy health state"],
               ["Dashboard latency",     "p95 < 200 ms",            "Operator workflow speed"],
               ["Silent gap detection",  "< 5 min to page",         "No stale green dashboards"],
               ["Accepted event loss",   "0 after durable bus",     "Replayable correctness"],
           ],
           col_widths=[2.45, 2.0, 7.45])

    _tb(s, 0.85, 5.85, 11.5, 0.45,
        f"{ICONS['diamond']}  Guardrail: every MVP feature must improve customer "
        "actionability or de-risk the platform launch.", 14, True, C["ink"])
    _footer(s, 2)

    _notes(s, """\
PRESENTER NOTES — Slide 2: Scope, SLOs & Guardrails
====================================================

PURPOSE
Show the audience you scope BEFORE designing. Most SEM candidates jump to
architecture; you start with requirements discipline.

KEY TALKING POINTS
1. MVP limits to 4 anomaly types — high signal, customer-actionable.
   Everything else is explicitly deferred to reduce delivery risk.
2. SLO table: call out "silent gap detection" — this is the most dangerous
   failure mode in a security product (stale green dashboards).
3. GA scope adds alert lifecycle, suppression, and shadow migration.
   Shadow validation is critical: we do NOT flip source-of-truth until we can
   quantify false-positive/negative rates.
4. Guardrail: prevents scope creep. Every feature request during MVP must pass
   the "actionability or de-risk" test.

DATA POINTS TO MENTION
• Billions of events/day at millions of agents.
• Multi-tenant isolation with fair-usage controls.
• At-least-once delivery with idempotent processing.

ANTICIPATED QUESTIONS
• "Why not add notifications in MVP?"
  → Notifications add blast-radius risk; we need correctness first.
• "How do you measure the 5-minute freshness SLO?"
  → Synthetic events flowing end-to-end (preview Slide 6).

TRANSITION → "With scope locked, let me walk you through the reference
architecture that satisfies these SLOs."
""")


def slide_03_architecture(prs):
    """Reference Architecture"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 3, f"{ICONS['gear']}  Reference Architecture",
           "Cloud-neutral, replayable, and intentionally separated into hot, read, and cold paths")

    # Main pipeline boxes
    stages = [
        ("Agents",             C["soft_teal"],   0.45),
        ("Ingestion\nGateway", C["soft_blue"],   1.9),
        ("Durable\nEvent Bus", C["soft_violet"],  3.55),
        ("Stream\nProcessing", C["soft_green"],   5.2),
        ("Rules +\nCoalescing",C["soft_amber"],  6.85),
        ("Alert\nState Store", C["soft_teal"],   8.5),
        ("Read Model\n+ Search",C["soft_blue"],  10.15),
        ("APIs\n+ Console",    C["soft_violet"], 11.8),
    ]
    y = 2.0
    for i, (label, fill, x) in enumerate(stages):
        _box(s, x, y, 1.25, 0.85, label, fill, C["line"], size=9.6)
        if i < len(stages) - 1:
            _arrow(s, x + 1.28, y + 0.42, stages[i + 1][2] - 0.05, y + 0.42, C["slate"], 1.2)

    # Data Lake (cold path)
    _box(s, 3.55, 4.2, 1.25, 0.72, "Data Lake\nAudit / Replay", C["soft_green"], C["line"], size=9.2)
    _arrow(s, 4.18, 2.88, 4.18, 4.15, C["green"], 1.2)

    # Lifecycle Store
    _box(s, 8.5, 4.2, 1.25, 0.72, "Lifecycle\nStore", C["soft_amber"], C["line"], size=9.2)
    _arrow(s, 9.12, 2.88, 9.12, 4.15, C["amber"], 1.2)

    # Technology annotations
    _tb(s, 1.85, 3.08, 1.35, 0.22, "Kafka / PubSub", 7, False, C["muted"], PP_ALIGN.CENTER)
    _tb(s, 3.5,  3.08, 1.35, 0.22, "Kafka / Kinesis", 7, False, C["muted"], PP_ALIGN.CENTER)
    _tb(s, 5.15, 3.08, 1.35, 0.22, "Flink / KStreams", 7, False, C["muted"], PP_ALIGN.CENTER)
    _tb(s, 8.45, 3.08, 1.35, 0.22, "DDB / Cassandra", 7, False, C["muted"], PP_ALIGN.CENTER)
    _tb(s, 10.1, 3.08, 1.35, 0.22, "OpenSearch / ES", 7, False, C["muted"], PP_ALIGN.CENTER)

    # Path summary cards
    _card(s, 0.75, 5.45, 3.5, 0.95,
          f"{ICONS['bolt']}  Hot Path",
          "Event-time detection, coalescing, alert transitions, freshness SLO.", C["teal"])
    _card(s, 4.9,  5.45, 3.5, 0.95,
          f"{ICONS['chart']}  Read Path",
          "Precomputed dashboard views, bounded queries, freshness metadata.", C["blue"])
    _card(s, 9.05, 5.45, 3.5, 0.95,
          f"{ICONS['lock']}  Cold Path",
          "Long-term retention, backfill, analytics, and rule tuning.", C["green"])
    _footer(s, 3)

    _notes(s, """\
PRESENTER NOTES — Slide 3: Reference Architecture
===================================================

PURPOSE
Present a cloud-neutral, replayable architecture. Emphasize the
SEPARATION OF CONCERNS (hot / read / cold) — this is what makes the
system scalable and operable.

WALK-THROUGH (left to right, ~90 seconds)
1. Agents → Ingestion Gateway: authentication, rate-limiting, routing.
   Gateway team owns this; we consume via a versioned telemetry contract.
2. Durable Event Bus: partitioned by tenant + agent ID.
   At-least-once delivery; replay is the key capability.
3. Stream Processing → Rules + Coalescing: real-time anomaly detection.
   This is where the 5-min SLO is enforced.
4. Alert State Store: write-optimized, tenant/agent-keyed.
5. Read Model + Search: pre-aggregated views for console queries.
6. APIs + Console: bounded queries, pagination, freshness metadata.

COLD PATH
• Data Lake stores every event for audit, replay, and rule tuning.
• Lifecycle Store tracks alert history for compliance.

TECHNOLOGY CHOICES
• Deliberately cloud-neutral labels. In the interview, mention:
  Kafka or PubSub, Flink or Kafka Streams, DynamoDB or Cassandra,
  OpenSearch or Elasticsearch, S3/GCS for the data lake.
• Do NOT over-commit to specific vendors; show flexibility.

DESIGN PRINCIPLES
• Write path never shares resources with read path.
• Every stage emits metrics for pipeline-health monitoring (Slide 6).
• Replay-from-offset is the primary recovery mechanism.

ANTICIPATED QUESTIONS
• "Why not a single database?"
  → Different access patterns require different stores; forcing one
    store to serve writes + queries + analytics creates contention.
• "How do you handle schema evolution?"
  → Versioned telemetry contract with backward-compat validation.

TRANSITION → "Let me zoom into the read path and how we hit the 200ms
latency target."
""")


def slide_04_read_path(prs):
    """Read Path And Latency Budget"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 4, f"{ICONS['bolt']}  Read Path And Latency Budget",
           "The UI should query health state, not raw telemetry")

    # Read pipeline flow
    read_stages = [
        ("Console\nDashboard",       C["soft_violet"], 0.8),
        ("Health APIs\nBounded Queries", C["soft_blue"], 3.35),
        ("Read Model\nAggregates",   C["soft_teal"],   5.9),
        ("Search Index\nFilters",    C["soft_green"],  8.45),
        ("Alert\nState",             C["soft_amber"],  10.95),
    ]
    for i, (label, fill, x) in enumerate(read_stages):
        w = 2.0 if i < 4 else 1.6
        _box(s, x, 1.45, w, 0.75, label, fill, C["line"], size=9.6)
    for x1, x2 in [(2.8, 3.35), (5.35, 5.9), (7.9, 8.45), (10.45, 10.95)]:
        _arrow(s, x1, 1.83, x2, 1.83)

    # Latency budget table
    _table(s, 0.8, 3.0, 7.4, 2.3,
           [f"{ICONS['clock']}  Budget Segment", "Target", "Design Choice"],
           [
               ["Auth / gateway",         "20–40 ms",  "Tenant-scoped claims & request shaping"],
               ["API orchestration",      "40–60 ms",  "No fanout over raw telemetry"],
               ["Read / search",          "70–90 ms",  "Precomputed views & bounded filters"],
               ["Serialization / network","20–30 ms",  "Slim response models & pagination"],
           ],
           col_widths=[2.1, 1.45, 3.85])

    # Design-principle cards
    _card(s, 8.75, 3.05, 3.75, 0.9,
          f"{ICONS['warning']}  Graceful Degradation",
          "If search lags, serve current alert state with explicit freshness metadata.", C["amber"])
    _card(s, 8.75, 4.15, 3.75, 0.9,
          f"{ICONS['lock']}  Query Discipline",
          "Pagination, capped time windows, no tenant-wide scans from UI paths.", C["blue"])
    _card(s, 8.75, 5.25, 3.75, 0.9,
          f"{ICONS['check']}  Correctness Model",
          "At-least-once ingestion plus idempotent alert transitions.", C["green"])
    _footer(s, 4)

    _notes(s, """\
PRESENTER NOTES — Slide 4: Read Path & Latency Budget
=======================================================

PURPOSE
Demonstrate that the 200ms p95 target is achievable through design, not
hope. Show you've thought through every millisecond of the request path.

WALK-THROUGH
1. Console → APIs → Read Model → Search → Alert State.
   Data flows RIGHT to LEFT in the query path.
2. Latency budget adds up to ~150–220 ms.  We have headroom under p95.
3. Key insight: the UI queries PRECOMPUTED VIEWS, not raw events.

IMPORTANT DESIGN CHOICES
• Graceful Degradation: if search index lags, the API returns current
  alert state with freshness metadata.  The console shows "data may be
  delayed" instead of stale green.  This is critical for a security
  product.
• Query Discipline: bounded queries prevent runaway scans.  No
  "select * from events where tenant=X" from the dashboard.
• Correctness: at-least-once + idempotent transitions means we can
  safely replay without creating duplicate alerts.

LATENCY BUDGET DEEP-DIVE (if asked)
• Auth/gateway: JWT validation, tenant-scoped claims, rate shaping.
• Orchestration: single API call to the read model, no service fanout.
• Read store: OpenSearch query with pre-aggregated alert-count views.
• Serialization: protobuf or slim JSON, paginated to max 100 items.

ANTICIPATED QUESTIONS
• "What if a tenant has 100K agents?"
  → Pagination + server-side aggregation.  Dashboard shows top-N with
    drill-down, not a flat list.
• "How do you cache?"
  → Tenant-level dashboard summaries with TTL-based invalidation.
    Per-agent pages are not cached.

TRANSITION → "Now let me show you the hardest technical problem: alert
coalescing within the 5-minute SLO."
""")


def slide_05_coalescing(prs):
    """Alert Coalescing Design"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 5, f"{ICONS['fire']}  Alert Coalescing Design",
           "Turn noisy telemetry into actionable health state within the 5-minute SLO")

    _tb(s, 0.75, 1.25, 5.2, 0.28,
        "Example: 50 anti-tamper disabled events over 10 minutes", 12.5, True)

    # Timeline
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.85), Inches(2.05), Inches(5.7), Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb = C["line"]; ln.line.color.rgb = C["line"]
    event_xs = [0.95, 1.4, 1.8, 2.3, 2.9, 3.45, 4.0, 4.65, 5.2, 5.8, 6.25]
    for i, x in enumerate(event_xs):
        fill = C["rose"] if i <= 5 else C["amber"]
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(1.93), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = fill; dot.line.color.rgb = fill
    _pill(s, 0.85, 2.42, 1.2, 0.34, "t = 0",   C["soft_rose"], size=9)
    _pill(s, 5.55, 2.42, 1.1, 0.34, "t = 10 m", C["soft_amber"], size=9)

    # Coalescing pipeline
    _box(s, 7.0,  1.2, 1.55, 0.64, f"{ICONS['link']}  Key By\nTenant+Agent", C["soft_blue"], C["line"], size=9)
    _box(s, 9.0,  1.2, 1.55, 0.64, f"{ICONS['clock']}  Window\nState",      C["soft_teal"], C["line"], size=9)
    _box(s, 11.0, 1.2, 1.55, 0.64, f"{ICONS['check']}  Alert\nTransition",  C["soft_green"], C["line"], size=9)
    _arrow(s, 8.55, 1.52, 9.0, 1.52)
    _arrow(s, 10.55, 1.52, 11.0, 1.52)

    # State-machine / auditability / decision cards
    _card(s, 0.85, 3.35, 3.6, 1.05,
          f"{ICONS['gear']}  State Machine",
          "Open → update evidence count → suppress duplicate → resolve when healthy signal appears.",
          C["teal"])
    _card(s, 4.75, 3.35, 3.6, 1.05,
          f"{ICONS['lock']}  Auditability",
          "Store raw event references and rule version, not every raw event in the UI path.",
          C["blue"])
    _card(s, 8.65, 3.35, 3.6, 1.05,
          f"{ICONS['target']}  Decision",
          "Use stream processing for coalescing; keep rule scope narrow for MVP.",
          C["green"])

    _bullets(s, 0.95, 5.05, 11.2, 0.95, [
        f"{ICONS['warning']}  Why not cron-first: late detection, expensive scans, duplicate intermediate state, and weak replay semantics.",
        f"{ICONS['check']}  How to control complexity: fixed MVP rules, explicit windows, idempotent updates, per-rule metrics, and rollback flags.",
    ], 11)
    _footer(s, 5)

    _notes(s, """\
PRESENTER NOTES — Slide 5: Alert Coalescing Design
====================================================

PURPOSE
This is the HARDEST TECHNICAL PROBLEM in the design.  Show that you
understand the trade-off and can make an evidence-based decision.

EXAMPLE WALK-THROUGH (point to timeline)
• 50 anti-tamper events from one agent over 10 minutes.
• Without coalescing, the dashboard shows 50 separate alerts → noise.
• With coalescing: key by tenant+agent → tumbling window → state machine
  emits ONE alert with evidence count.

STATE MACHINE TRANSITIONS
  Open → Update Evidence → Suppress Duplicate → Resolve (healthy signal)
  Each transition is idempotent using (tenant, agent, rule, window) key.

WHY STREAM PROCESSING OVER CRON
1. Freshness: the 5-min SLO cannot be reliably met with periodic batch.
2. Scale: scanning billions of events in cron is expensive.
3. Replay: stream offsets give free replay; cron needs custom backfill.
4. Intermediate state: cron creates half-processed records.

COMPLEXITY CONTROLS (critical for credibility)
• Only 4 rules in MVP — no generic rule engine.
• Explicit tumbling windows (configurable per rule).
• Idempotent state updates.
• Per-rule metrics: processing lag, emit rate, error rate.
• Rollback flags: disable individual rules without redeploying.

THIS IS THE ENGINEER A vs ENGINEER B DEBATE
• Preview that you expect stream processing to win, but you'll run a
  structured decision process (Slide 9).

ANTICIPATED QUESTIONS
• "What if a rule is buggy?"
  → Feature flag per rule.  Disable instantly.  Roll back to last-known-
    good rule version.
• "How do you test coalescing logic?"
  → Unit tests with synthetic event sequences.  Integration tests with
    replayed production data (scrubbed).

TRANSITION → "How do we know all of this is working correctly in
production?  Let me walk through the reliability model."
""")


def slide_06_reliability(prs):
    """Reliability Model And Sev-1 Runbook"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 6, f"{ICONS['shield']}  Reliability Model And Sev-1 Runbook",
           "Design for silent failure detection, not just service uptime")

    # Observability pipeline stages
    obs_stages = [
        ("Ingress\nrate",     0.8,  C["soft_blue"]),
        ("Bus\nlag",          2.5,  C["soft_violet"]),
        ("Processor\nlag",    4.2,  C["soft_teal"]),
        ("Store\nwrites",     5.9,  C["soft_green"]),
        ("Index\nfreshness",  7.6,  C["soft_amber"]),
        ("Synthetic\nprobe",  9.3,  C["soft_rose"]),
        ("API\np95",          11.0, C["soft_blue"]),
    ]
    for i, (label, x, fill) in enumerate(obs_stages):
        _box(s, x, 1.35, 1.15, 0.72, label, fill, C["line"], size=8.8)
        if i < len(obs_stages) - 1:
            _arrow(s, x + 1.15, 1.71, obs_stages[i + 1][1] - 0.05, 1.71)

    # Sev-1 runbook table
    _table(s, 0.75, 2.75, 11.85, 2.8,
           [f"{ICONS['fire']}  Incident Step", "SEM / IC Focus", "Technical Focus"],
           [
               ["Declare",   "Single IC, comms owner, customer impact clock", "Freeze risky deploys"],
               ["Scope",     "Tenant, region, event type, time window",       "Compare stage metrics"],
               ["Mitigate",  "Choose customer-safe degraded mode",            "Failover, pause rules, restart consumers"],
               ["Recover",   "Track backfill completion",                      "Replay from durable offsets"],
               ["Prevent",   "Postmortem and owners",                          "Synthetic gaps, gates, runbooks"],
           ],
           col_widths=[2.0, 4.35, 5.5])

    _tb(s, 0.9, 6.05, 11.2, 0.35,
        f"{ICONS['warning']}  Operating principle: a stale green Health Center is worse "
        "than an explicit degraded state.", 14, True, C["rose"])
    _footer(s, 6)

    _notes(s, """\
PRESENTER NOTES — Slide 6: Reliability & Sev-1 Runbook
========================================================

PURPOSE
This slide is CRITICAL. It demonstrates operational maturity — the #1
differentiator for a Senior EM candidate.

OBSERVABILITY PIPELINE (top diagram)
Walk through each stage LEFT to RIGHT:
1. Ingress rate: events accepted by the gateway per second.
2. Bus lag: consumer offset lag on the durable event bus.
3. Processor lag: processing latency vs event time.
4. Store writes: write throughput to alert state store.
5. Index freshness: lag between alert-store writes and search-index
   availability.
6. Synthetic probe: end-to-end synthetic events per tenant cohort.
7. API p95: dashboard response time.

KEY INSIGHT: "A Kubernetes pod being UP is not proof that the pipeline
is processing events correctly."

SEV-1 RUNBOOK — WALK THROUGH EACH ROW
1. DECLARE: Single incident commander.  Start customer-impact clock.
   Freeze risky deployments.
2. SCOPE: Determine blast radius — which tenants, regions, event types,
   time window.  Compare metrics across stages.
3. MITIGATE: Choose a CUSTOMER-SAFE degraded mode.  Pause bad rules.
   Restart consumers.  Failover to standby.
4. RECOVER: Replay from durable offsets.  Backfill missing alerts.
   Track completion.
5. PREVENT: Blameless postmortem.  Add synthetic gaps.  Improve gates.

OPERATING PRINCIPLE (bottom, in red)
"A stale green Health Center is worse than an explicit degraded state."
This is the single most important sentence in the deck.  It shows you
understand security-product operations.

ANTICIPATED QUESTIONS
• "How long does replay take?"
  → Depends on gap duration.  For a 1-hour gap with billions of events,
    expect 20-30 minutes with parallelized consumers.
• "Who gets paged?"
  → On-call SRE + Health Center tech lead + SEM for customer comms.

TRANSITION → "Let me show you the multi-quarter roadmap with delivery
gates."
""")


def slide_07_roadmap(prs):
    """Roadmap: MVP, GA, Platformization"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 7, f"{ICONS['rocket']}  Roadmap: MVP, GA, Platformization",
           "Multi-quarter plan with customer-safety gates")

    # Quarter headers
    quarters = [
        (f"{ICONS['plant']}  Q1 MVP",           2.55),
        (f"{ICONS['rocket']}  Q2 GA",           5.55),
        (f"{ICONS['gear']}  Q3 Platform",       8.55),
        (f"{ICONS['trophy']}  Q4 Reduce Legacy",11.1),
    ]
    for label, x in quarters:
        _pill(s, x, 1.25, 1.55, 0.38, label, C["ink"], C["white"], 9.5)

    # Swim lanes
    lanes = [
        ("Telemetry +\nProcessing", C["teal"],
         ["Contract\nAdapter\nTop 4 rules", "Replay\nScale test\nIdempotency",
          "Reusable\ncoalescing\nframework", "Broader\nhealth catalog"]),
        ("Console +\nAPIs", C["blue"],
         ["Dashboard\nAlert list\nFeature flag", "Lifecycle\nSearch\nDrill-down",
          "Trends\nReports\nNotifications", "Legacy UI\nconsolidation"]),
        ("Operations", C["green"],
         ["SLOs\nSynthetic\nRunbooks", "Game days\nCanaries\nGA gates",
          "Cost model\nOwnership\nOnboarding", "Decom\nrunbooks"]),
        ("Legacy\nOffline", C["amber"],
         ["Stabilize\nfalse offline", "Shadow\ncompare",
          "Tenant\nmigration", "Retire\nsafe paths"]),
    ]
    y = 1.85
    for lane_name, accent, items in lanes:
        _tb(s, 0.65, y + 0.14, 1.65, 0.55, lane_name, 9.2, True, C["ink"])
        for i, item in enumerate(items):
            bw = 1.9 if i < 3 else 1.65
            _box(s, 2.55 + i * 3.0, y, bw, 0.7, item,
                 C["white"], accent, size=8.1)
            if i < 3:
                _arrow(s, 4.45 + i * 3.0, y + 0.35,
                       5.0 + i * 3.0, y + 0.35, C["line"], 1)
        y += 1.05

    # Gate cards
    _card(s, 0.75, 6.25, 3.8, 0.65,
          f"{ICONS['check']}  MVP Gate",
          "End-to-end synthetic, top 4 signals, SLO dashboards, allowlist rollout.", C["teal"])
    _card(s, 4.8,  6.25, 3.8, 0.65,
          f"{ICONS['target']}  GA Gate",
          "Scale, replay, correctness, shadow validation, incident readiness.", C["blue"])
    _card(s, 8.85, 6.25, 3.8, 0.65,
          f"{ICONS['warning']}  Migration Gate",
          "Measured FP/FN rate and rollback path per tenant cohort.", C["amber"])
    _footer(s, 7)

    _notes(s, """\
PRESENTER NOTES — Slide 7: Roadmap
====================================

PURPOSE
Show a credible, multi-quarter execution plan with EXPLICIT GATES that
prevent premature rollout.

SWIM-LANE WALK-THROUGH (~60 seconds)
• Telemetry + Processing: contract → adapt → top 4 rules (Q1) →
  scale/replay/idem (Q2) → reusable framework (Q3) → broader catalog (Q4).
• Console + APIs: dashboard/alerts/flags (Q1) → lifecycle/search (Q2) →
  trends/reports/notifs (Q3) → legacy UI consolidation (Q4).
• Operations: SLOs/synthetic/runbooks from DAY ONE (Q1) → game days/
  canaries (Q2) → cost model/onboarding (Q3) → decom runbooks (Q4).
• Legacy Offline: stabilize false-offline (Q1) → shadow compare (Q2) →
  tenant migration (Q3) → retire safe paths (Q4).

GATE DEFINITIONS (bottom cards)
1. MVP Gate: cannot GA without end-to-end synthetic events, SLO dashboards,
   and controlled allowlist rollout.
2. GA Gate: cannot go broad without scale testing, replay validation,
   shadow comparison, and incident readiness.
3. Migration Gate: cannot retire legacy without measured FP/FN rates and
   a per-tenant rollback path.

KEY MESSAGE: "Operations is a first-class swim lane, not an afterthought.
We ship SLOs and runbooks before we ship features."

ANTICIPATED QUESTIONS
• "Can you compress Q1 and Q2?"
  → Possibly, if gateway dependency resolves early.  But I'd rather ship
    a bulletproof MVP than a rushed GA.
• "When is legacy fully retired?"
  → Q4 for safe paths; some long-tail tenants may extend into Q5.
    Explicit tenant-by-tenant migration with rollback.

TRANSITION → "Let me address the two biggest disruptions: the gateway
delay and the legacy escalation spike."
""")


def slide_08_disruption(prs):
    """Dependency Delay And Operational Drain"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 8, f"{ICONS['warning']}  Dependency Delay And Operational Drain",
           "Absorb real-world disruption without pretending capacity is infinite")

    _card(s, 0.75, 1.25, 5.8, 1.0,
          f"{ICONS['link']}  Gateway Delayed By Two Months",
          "Freeze telemetry contract, build adapter boundary, generate synthetic/"
          "replayed streams, negotiate thin-slice routing for highest-value signals.", C["violet"])
    _card(s, 6.95, 1.25, 5.6, 1.0,
          f"{ICONS['fire']}  Legacy Offline Escalations +40%",
          "Treat as customer trust risk; stabilize with a short-lived lane, not a "
          "full-team roadmap derailment.", C["rose"])

    # Team allocation stacked bar (org chart / allocation visual)
    _tb(s, 0.85, 2.85, 5.0, 0.28,
        f"{ICONS['people']}  10-person team allocation during escalation", 12.5, True)
    segments = [
        ("Health Center build", 6, C["blue"]),
        ("Legacy stabilization", 2, C["rose"]),
        ("QA / release",        1, C["green"]),
        ("Tech lead coordination", 1, C["amber"]),
    ]
    x = 0.85
    total_w = 7.2
    for label, count, color in segments:
        w = total_w * count / 10
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(3.35), Inches(w), Inches(0.55))
        rect.fill.solid(); rect.fill.fore_color.rgb = color
        rect.line.color.rgb = C["white"]
        # Count label inside bar
        _tb(s, x + 0.05, 3.42, w - 0.1, 0.22, str(count), 10, True, C["white"], PP_ALIGN.CENTER)
        # Label below bar
        _tb(s, x, 4.05, w, 0.38, label, 8.2, False, C["muted"], PP_ALIGN.CENTER)
        x += w

    # Strategy cards
    _card(s, 8.45, 2.85, 4.0, 0.85,
          f"{ICONS['check']}  Exit Criteria",
          "Legacy lane ends after top root cause patch, instrumentation, and "
          "support-volume trend review.", C["green"])
    _card(s, 8.45, 3.95, 4.0, 0.85,
          f"{ICONS['warning']}  Escalation Rule",
          "Increase allocation only for active Sev-1/Sev-2 impact across major tenants.", C["amber"])
    _card(s, 8.45, 5.05, 4.0, 0.85,
          f"{ICONS['gear']}  Learning Loop",
          "Each false-offline root cause becomes a Health Center test, rule "
          "refinement, or migration guardrail.", C["blue"])
    _footer(s, 8)

    _notes(s, """\
PRESENTER NOTES — Slide 8: Disruption Management
==================================================

PURPOSE
This is where you demonstrate EXECUTION LEADERSHIP.  Two real-world
disruptions that most candidates handle poorly:
1. A dependency team is delayed.
2. A legacy system is burning.

GATEWAY DELAY RESPONSE
• Do NOT idle the team.
• Freeze a versioned telemetry contract.
• Build an adapter boundary against that contract.
• Generate synthetic + replayed streams.
• Unblock: processors, schemas, APIs, UI, QA, observability, load tests.
• Negotiate thin-slice integration: route 1-2 high-value signals first.
• Only dependency-bound scope shifts; everything else continues.

LEGACY ESCALATION RESPONSE
• Treat as customer trust risk, NOT background maintenance.
• Create a SHORT-LIVED stabilization lane: 2 engineers + QA support.
• Keep 6 engineers on Health Center MVP.  1 tech lead coordinates.
• Mission: instrument → root-cause → patch → reduce support volume.
• Exit criteria are explicit: lane closes after patch + instrumentation +
  support-volume trend review.

TEAM ALLOCATION BAR (visual)
• Walk through: 6 build / 2 legacy / 1 QA / 1 coordination = 10 total.
• Emphasize: this is a TEMPORARY allocation, not a permanent split.
• Escalation rule: only increase legacy if Sev-1/Sev-2 across major tenants.

LEARNING LOOP
• Every false-offline root cause feeds back into Health Center:
  becomes a test case, a rule refinement, or a migration guardrail.
• This turns the legacy pain into strategic value.

ANTICIPATED QUESTIONS
• "What if the gateway delay extends to 4 months?"
  → We can ship MVP entirely on synthetic + existing heartbeat streams.
    Production gateway integration becomes a GA gate.
• "What if escalations consume the whole team?"
  → The escalation rule is explicit.  SEM protects roadmap unless
    active Sev-1/Sev-2 across major tenants.

TRANSITION → "Now let me address the people challenge: resolving the
technical conflict between the two senior engineers."
""")


def slide_09_conflict(prs):
    """Technical Conflict Resolution"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 9, f"{ICONS['people']}  Technical Conflict Resolution",
           "Use evidence, decision criteria, and ownership to align senior engineers")

    # Org chart: Engineer A ←→ SEM ←→ Engineer B
    _box(s, 0.85,  1.35, 2.25, 0.78,
         f"{ICONS['bolt']}  Engineer A\nStream processing", C["soft_teal"], C["teal"], size=9.3)
    _box(s, 10.25, 1.35, 2.25, 0.78,
         f"{ICONS['gear']}  Engineer B\nDB + cron",         C["soft_amber"], C["amber"], size=9.3)
    _box(s, 4.95,  1.25, 3.2, 0.98,
         f"{ICONS['people']}  SEM Role\nTurn disagreement into\na decision system",
         C["soft_blue"], C["blue"], size=10)
    _arrow(s, 3.1, 1.74, 4.9, 1.74)
    _arrow(s, 10.25, 1.74, 8.2, 1.74)

    # 4-step resolution process
    steps = [
        (f"{ICONS['target']}  1 Criteria",  "Freshness, scale, cost, operability, delivery risk"),
        (f"{ICONS['gear']}  2 Spike",       "One week with realistic cardinality and failure cases"),
        (f"{ICONS['check']}  3 Decision",   "ADR with assumptions, rollback, and revisit trigger"),
        (f"{ICONS['people']}  4 Commit",    "Both engineers own part of the winning design"),
    ]
    accents = [C["teal"], C["blue"], C["violet"], C["green"]]
    for i, (h, b) in enumerate(steps):
        _card(s, 0.9 + i * 3.05, 3.0, 2.55, 1.1, h, b, accents[i])
        if i < 3:
            _arrow(s, 0.9 + i * 3.05 + 2.55, 3.55,
                   0.9 + (i + 1) * 3.05 - 0.12, 3.55, C["line"], 1.2)

    _tb(s, 0.9, 5.05, 11.5, 0.45,
        f"{ICONS['diamond']}  Expected decision: stream processing for coalescing, "
        "constrained to MVP rules and strong operability controls.", 14, True, C["ink"])
    _bullets(s, 1.05, 5.65, 10.8, 0.8, [
        f"{ICONS['bolt']}  Engineer A leads processing design and SLO instrumentation.",
        f"{ICONS['gear']}  Engineer B leads simplicity controls: bounded scope, failure-mode "
        "review, cost model, and rollback plan.",
    ], 10.8)
    _footer(s, 9)

    _notes(s, """\
PRESENTER NOTES — Slide 9: Conflict Resolution
================================================

PURPOSE
This is the PEOPLE LEADERSHIP test.  The interviewer wants to see that
you can lead senior engineers through disagreement WITHOUT:
• Pulling rank ("just do it my way")
• Avoiding the decision ("let's try both")
• Making it personal ("A is right, B is wrong")

ORG CHART VISUAL (top)
• Engineer A: advocates stream processing (Kafka Streams / Flink).
• Engineer B: advocates DB writes + cron aggregation.
• SEM (you): turns disagreement into a decision system.

4-STEP PROCESS
1. CRITERIA: Define decision criteria BEFORE debating solutions.
   Freshness, scale, cost, operability, delivery risk.
2. SPIKE: One week with realistic data.  Not a toy prototype.
   Realistic cardinality (millions of agents), duplicate-event rates,
   and failure cases (what happens when a processor crashes?).
3. DECISION: Architecture Decision Record (ADR).
   Document: assumptions, constraints, decision, trade-offs, rollback
   trigger, revisit conditions.
4. COMMIT: Both engineers own part of the winning design.
   Engineer A leads the processing layer.
   Engineer B leads simplicity controls and rollback planning.
   This turns disagreement into a STRONGER DESIGN.

WHY STREAM PROCESSING WINS
• 5-minute freshness SLO at billions of events/day.
• Stream-time windows naturally model coalescing.
• Durable offsets enable replay.
• Cron risks: late detection, expensive scans, duplicate state.

BUT with constraints:
• Only 4 rules in MVP.
• No generic rule engine.
• Feature flags per rule.
• Strong operability: metrics, dashboards, runbooks, rollback.

ANTICIPATED QUESTIONS
• "What if they still disagree after the spike?"
  → The SEM makes the call.  The ADR documents it.  Both engineers
    sign off on the document, not just the decision.
• "Have you done this before?"
  → [Prepare a real example from your experience]

TRANSITION → "Let me close with the major risks, mitigations, and the
outcomes we expect to deliver."
""")


def slide_10_risks(prs):
    """Major Risks, Mitigations, And Outcomes"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _title(s, 10, f"{ICONS['trophy']}  Major Risks, Mitigations, And Outcomes",
           "Close with operating maturity, not a technology shopping list")

    # ── Risk matrix 4×4 ──
    _tb(s, 0.75, 1.22, 4.6, 0.28, f"{ICONS['warning']}  Risk Matrix", 12.5, True)
    x0, y0, cell = 0.8, 1.65, 0.82
    # Severity coloring
    severity_colors = [C["soft_green"], C["soft_amber"], C["soft_rose"], C["soft_rose"]]
    for i in range(4):
        for j in range(4):
            color = severity_colors[min(i + j, 3)]
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x0 + i * cell),
                                       Inches(y0 + (3 - j) * cell),
                                       Inches(cell), Inches(cell))
            rect.fill.solid(); rect.fill.fore_color.rgb = color
            rect.line.color.rgb = C["white"]
    # Axis labels
    _tb(s, 0.8, 5.0, 3.3, 0.25, f"Likelihood {ICONS['arrow_r']}", 8.5, False, C["muted"], PP_ALIGN.CENTER)
    _tb(s, 0.25, 2.8, 0.45, 0.5, "Impact", 8.5, False, C["muted"])

    # Risk dots
    risks = [
        ("G", 3, 3, C["violet"]),  # Gateway delay
        ("S", 2, 3, C["blue"]),    # Stream complexity
        ("L", 3, 2, C["rose"]),    # Legacy drain
        ("Q", 2, 2, C["amber"]),   # Query latency
        ("N", 1, 2, C["green"]),   # Noise / alert fatigue
    ]
    for label, ix, iy, color in risks:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x0 + ix * cell + 0.22),
                                  Inches(y0 + (3 - iy) * cell + 0.2),
                                  Inches(0.36), Inches(0.36))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.color.rgb = color
        _tb(s, x0 + ix * cell + 0.27, y0 + (3 - iy) * cell + 0.27,
            0.2, 0.12, label, 7, True, C["white"], PP_ALIGN.CENTER)

    # Risk detail cards
    _card(s, 4.7,  1.35, 3.7, 0.75,
          f"{ICONS['link']}  G  Gateway Delay",
          "Adapter contract, synthetic streams, thin-slice integration.", C["violet"])
    _card(s, 8.75, 1.35, 3.7, 0.75,
          f"{ICONS['gear']}  S  Stream Complexity",
          "Narrow MVP rules, SLOs, runbooks, rollback flags.", C["blue"])
    _card(s, 4.7,  2.35, 3.7, 0.75,
          f"{ICONS['fire']}  L  Legacy Drain",
          "Two-engineer stabilization lane with exit criteria.", C["rose"])
    _card(s, 8.75, 2.35, 3.7, 0.75,
          f"{ICONS['clock']}  Q  Query Latency",
          "Precomputed read models, bounded APIs, freshness metadata.", C["amber"])
    _card(s, 4.7,  3.35, 3.7, 0.75,
          f"{ICONS['warning']}  N  Noise / Alert Fatigue",
          "Coalescing, suppression, tenant rollout, support feedback loop.", C["green"])

    # Expected business outcomes
    _tb(s, 4.75, 4.75, 7.7, 0.35,
        f"{ICONS['star']}  Expected Business Outcomes", 13, True, C["ink"])
    _bullets(s, 4.9, 5.22, 7.4, 1.05, [
        f"{ICONS['check']}  More trustworthy agent-health visibility for customers and support.",
        f"{ICONS['chart']}  Reduced false-offline escalations and less legacy operational drag.",
        f"{ICONS['rocket']}  Reusable Agent Platform foundation for future health and alerting use cases.",
    ], 11)
    _footer(s, 10)

    _notes(s, """\
PRESENTER NOTES — Slide 10: Risks, Mitigations & Outcomes
===========================================================

PURPOSE
Close strong. This slide should leave the interviewer confident that you
can OPERATE, not just design.

RISK MATRIX WALK-THROUGH
Point to each dot and explain:

G (Gateway Delay) — HIGH likelihood, HIGH impact.
  Mitigation: adapter contract, synthetic streams, thin-slice integration.
  Why it's manageable: we decouple downstream development from the
  gateway timeline.

S (Stream Complexity) — MEDIUM likelihood, HIGH impact.
  Mitigation: narrow MVP rules, SLOs, runbooks, rollback flags.
  Why it's manageable: complexity is bounded by scope control.

L (Legacy Drain) — HIGH likelihood, MEDIUM impact.
  Mitigation: dedicated stabilization lane with explicit exit criteria.
  Why it's manageable: short-lived, bounded allocation.

Q (Query Latency) — MEDIUM likelihood, MEDIUM impact.
  Mitigation: precomputed read models, bounded APIs, freshness metadata.
  Why it's manageable: standard read-path optimization.

N (Noise / Alert Fatigue) — LOW likelihood, MEDIUM impact.
  Mitigation: coalescing, suppression, tenant rollout, feedback loop.
  Why it's manageable: tenant-by-tenant rollout catches noise early.

EXPECTED OUTCOMES (bottom)
1. More trustworthy agent-health visibility.
2. Reduced false-offline escalations → less support burden.
3. Reusable Agent Platform foundation → future leverage.

CLOSING STATEMENT (memorize this)
"The Health Center program succeeds if it improves customer trust while
creating a durable Agent Platform foundation. The architecture handles
hyperscale telemetry, but the bigger test is execution judgment: narrow
MVP scope, protect customer-facing reliability, absorb dependency delays,
resolve technical conflict with evidence, and keep the team engaged
through operational pressure."

ANTICIPATED QUESTIONS
• "What would you cut if you had only one quarter?"
  → Ship top 2 signals (offline + anti-tamper) with synthetic SLOs on
    allowlisted tenants.  No lifecycle, no migration, no search.
• "How do you measure success?"
  → 50% reduction in false-offline escalations within 2 quarters.
    SLO dashboards showing 99%+ freshness compliance.
• "What's the biggest lesson from a past project?"
  → [Prepare a real example]
""")


# ──────────────────────────────────────────────────────────────
# Build & convert
# ──────────────────────────────────────────────────────────────

def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for fn in [
        slide_01_title,
        slide_02_scope,
        slide_03_architecture,
        slide_04_read_path,
        slide_05_coalescing,
        slide_06_reliability,
        slide_07_roadmap,
        slide_08_disruption,
        slide_09_conflict,
        slide_10_risks,
    ]:
        fn(prs)

    prs.save(str(PPTX_OUT))
    print(f"[OK] PPTX saved → {PPTX_OUT}")
    return PPTX_OUT


def convert_to_pdf(pptx_path: Path):
    """Convert PPTX to PDF using LibreOffice."""
    try:
        subprocess.run(
            [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(pptx_path.parent),
                str(pptx_path),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        print(f"[OK] PDF  saved → {PDF_OUT}")
    except FileNotFoundError:
        print("[WARN] LibreOffice not found. Install with:")
        print("       sudo apt-get install -y libreoffice-impress")
        print("       Then re-run this script.")
        return False
    except subprocess.CalledProcessError as e:
        print(f"[ERROR] LibreOffice conversion failed: {e.stderr}")
        return False
    return True


def main():
    pptx_path = build_pptx()
    convert_to_pdf(pptx_path)


if __name__ == "__main__":
    main()
